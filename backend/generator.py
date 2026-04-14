from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
import base64

def create_presentation_from_images(images_base64: list[str]) -> BytesIO:
    prs = Presentation()
    
    # Standard 16:9 widescreen dimensions (e.g., 10x5.625 inches or generic 13.333x7.5)
    # python-pptx default slide width/height: 10 inches by 7.5 inches for 4:3.
    # Widescreen 16:9 is typically 13.333 x 7.5 inches.
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_slide_layout = prs.slide_layouts[6]
    
    for idx, b64_str in enumerate(images_base64):
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Remove data:image/png;base64, prefix if it exists
        if "base64," in b64_str:
            b64_str = b64_str.split("base64,")[1]
            
        image_data = base64.b64decode(b64_str)
        img_stream = BytesIO(image_data)
        
        # Add picture filling the slide
        slide.shapes.add_picture(img_stream, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)

    ppt_stream = BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)
    
    return ppt_stream
